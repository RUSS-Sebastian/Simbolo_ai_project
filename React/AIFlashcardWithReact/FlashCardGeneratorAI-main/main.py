import os
import json
from flask import Flask, render_template, request, jsonify, send_file
from pypdf import PdfReader
from pptx import Presentation
import google.generativeai as genai

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Configure API Key
api_key = "AIzaSyAnt8NVqVam3PHiZVkuczGc490oHiXd-M0"
genai.configure(api_key=api_key)
model_gen = genai.GenerativeModel("gemini-1.5-flash")

# Functions to extract text from files
def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    return " ".join(page.extract_text() for page in reader.pages)

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    return "\n".join(
        shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
    )

def extract_text(file_path):
    if file_path.lower().endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        return ""

def generate_questions(text, number_of_cards, difficulty):
    prompt = (
        f"Generate exactly {number_of_cards} {difficulty} questions in a JSON array. "
        "Each item should have 'question' and 'answer'. Example: [{'question': '...', 'answer': '...'}] "
        f"Here's the text: {text}"
    )
    response = model_gen.generate_content(prompt)
    raw_result = response.text if hasattr(response, 'text') else response

    try:
        start_index = raw_result.index("[")
        end_index = raw_result.rindex("]") + 1
        cleaned_result = raw_result[start_index:end_index]
        questions = json.loads(cleaned_result)
    except (ValueError, json.JSONDecodeError):
        questions = []
    return questions

@app.route("/")
def index():
    return render_template('index.html')

@app.route("/generate", methods=["POST"])
def generate():
    file = request.files["file"]
    difficulty = request.form["difficulty"]
    num_questions = int(request.form["numQuestions"])

    if file:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(file_path)
        text = extract_text(file_path)

        if text:
            questions = generate_questions(text, num_questions, difficulty)
            os.remove(file_path)  # Clean up uploaded file
            if questions:
                return jsonify({"questions": questions})
            else:
                return jsonify({"error": "Failed to generate questions."})
        else:
            return jsonify({"error": "Failed to extract text from file."})
    return jsonify({"error": "No file uploaded."})

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5001))  # Default to 5001 if PORT is not set
    app.run(host="0.0.0.0", port=port, debug=True)
